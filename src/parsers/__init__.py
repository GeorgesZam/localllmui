"""
Document Parser Factory.

Implements Factory pattern for creating document parsers.
Allows easy registration of new parser types.
"""

import os
from typing import Dict, Callable, Optional, Any
from abc import ABC, abstractmethod

from patterns import Factory, FactoryItem


class DocumentParser(ABC):
    """
    Abstract base class for document parsers.

    All document parsers must inherit from this class
    and implement the parse method.
    """

    @abstractmethod
    def parse(self, file_path: str) -> str:
        """
        Parse a document and extract text.

        Args:
            file_path: Path to the document file

        Returns:
            Extracted text content
        """
        pass

    @abstractmethod
    def supports_file(self, file_path: str) -> bool:
        """
        Check if this parser supports the given file.

        Args:
            file_path: Path to check

        Returns:
            True if file is supported
        """
        pass

    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Get metadata about the document.

        Args:
            file_path: Path to the document

        Returns:
            Dictionary of metadata
        """
        return {
            'filename': os.path.basename(file_path),
            'extension': os.path.splitext(file_path)[1],
            'size': os.path.getsize(file_path) if os.path.exists(file_path) else 0
        }


class TextParser(DocumentParser):
    """Parser for plain text files."""

    TEXT_EXTENSIONS = {'.txt', '.md', '.csv', '.json', '.xml',
                       '.yaml', '.yml', '.html', '.htm', '.css',
                       '.py', '.js', '.ts', '.jsx', '.tsx', '.java',
                       '.c', '.cpp', '.h', '.hpp', '.cs', '.go', '.rs',
                       '.rb', '.php', '.sh', '.bash', '.zsh'}

    def parse(self, file_path: str) -> str:
        """Parse text file with encoding detection."""
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception:
                return ""
        return ""

    def supports_file(self, file_path: str) -> bool:
        """Check if file is a text file."""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.TEXT_EXTENSIONS


class PDFParser(DocumentParser):
    """Parser for PDF documents."""

    def __init__(self):
        self._has_pdf = self._check_pdf_support()

    def _check_pdf_support(self) -> bool:
        """Check if PyPDF2 is available."""
        try:
            import PyPDF2
            return True
        except ImportError:
            return False

    def parse(self, file_path: str) -> str:
        """Parse PDF document."""
        if not self._has_pdf:
            return ""

        try:
            import PyPDF2
            text_parts = []

            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)

                for page_num, page in enumerate(reader.pages):
                    try:
                        page_text = (page.extract_text() or "").strip()
                        if len(page_text) > 50:
                            text_parts.append(f"=== Page {page_num + 1} ===\n{page_text}")
                    except Exception:
                        continue

            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"[PDFParser] Error: {e}")
            return ""

    def supports_file(self, file_path: str) -> bool:
        """Check if file is a PDF."""
        return os.path.splitext(file_path)[1].lower() == '.pdf' and self._has_pdf


class DOCXParser(DocumentParser):
    """Parser for Word documents."""

    def __init__(self):
        self._has_docx = self._check_docx_support()

    def _check_docx_support(self) -> bool:
        """Check if python-docx is available."""
        try:
            from docx import Document
            return True
        except ImportError:
            return False

    def parse(self, file_path: str) -> str:
        """Parse DOCX document."""
        if not self._has_docx:
            return ""

        try:
            from docx import Document
            doc = Document(file_path)
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]

            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        text_parts.append(row_text)

            return "\n".join(text_parts)
        except Exception as e:
            print(f"[DOCXParser] Error: {e}")
            return ""

    def supports_file(self, file_path: str) -> bool:
        """Check if file is a DOCX."""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in {'.docx', '.doc'} and self._has_docx


class ExcelParser(DocumentParser):
    """Parser for Excel spreadsheets."""

    def __init__(self):
        self._has_excel = self._check_excel_support()

    def _check_excel_support(self) -> bool:
        """Check if openpyxl is available."""
        try:
            import openpyxl
            return True
        except ImportError:
            return False

    def parse(self, file_path: str) -> str:
        """Parse Excel spreadsheet."""
        if not self._has_excel:
            return ""

        try:
            import openpyxl
            text_parts = []
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) if c else "" for c in row)
                    if row_text.strip():
                        text_parts.append(row_text)

            wb.close()
            return "\n".join(text_parts)
        except Exception as e:
            print(f"[ExcelParser] Error: {e}")
            return ""

    def supports_file(self, file_path: str) -> bool:
        """Check if file is Excel."""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in {'.xlsx', '.xls'} and self._has_excel


class PPTXParser(DocumentParser):
    """Parser for PowerPoint presentations."""

    def __init__(self):
        self._has_pptx = self._check_pptx_support()

    def _check_pptx_support(self) -> bool:
        """Check if python-pptx is available."""
        try:
            from pptx import Presentation
            return True
        except ImportError:
            return False

    def parse(self, file_path: str) -> str:
        """Parse PowerPoint presentation."""
        if not self._has_pptx:
            return ""

        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = [s.text for s in slide.shapes
                             if hasattr(s, "text") and s.text]
                if slide_text:
                    text_parts.append(f"=== Slide {i} ===\n" + "\n".join(slide_text))

            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"[PPTXParser] Error: {e}")
            return ""

    def supports_file(self, file_path: str) -> bool:
        """Check if file is PowerPoint."""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in {'.pptx', '.ppt'} and self._has_pptx


class ImageParser(DocumentParser):
    """Parser for images using OCR."""

    def __init__(self, ocr_processor=None):
        self.ocr = ocr_processor

    def set_ocr_processor(self, ocr_processor) -> None:
        """Set the OCR processor."""
        self.ocr = ocr_processor

    def parse(self, file_path: str) -> str:
        """Parse image using OCR."""
        if not self.ocr or not self.ocr.available:
            return ""
        return self.ocr.ocr_image(file_path)

    def supports_file(self, file_path: str) -> bool:
        """Check if file is an image."""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in {'.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'}


class DocumentParserFactory:
    """
    Factory for creating document parsers.

    Implements Factory pattern with automatic parser selection
    based on file extension.
    """

    def __init__(self):
        self._parsers: Dict[str, DocumentParser] = {}
        self._factory = Factory()

        # Register default parsers
        self._register_defaults()

    def _register_defaults(self) -> None:
        """Register default document parsers."""
        self.register('text', TextParser, extensions=TextParser.TEXT_EXTENSIONS)
        self.register('pdf', PDFParser, extensions=['.pdf'])
        self.register('docx', DOCXParser, extensions=['.docx', '.doc'])
        self.register('excel', ExcelParser, extensions=['.xlsx', '.xls'])
        self.register('pptx', PPTXParser, extensions=['.pptx', '.ppt'])
        self.register('image', ImageParser, extensions=['.png', '.jpg', '.jpeg'])

    def register(self, name: str, parser_class: type,
                 extensions: list = None, set_as_default: bool = False) -> None:
        """
        Register a parser class.

        Args:
            name: Parser name
            parser_class: Parser class (must inherit from DocumentParser)
            extensions: List of file extensions this parser handles
            set_as_default: Whether to set as default parser
        """
        metadata = {'extensions': extensions or []}
        self._factory.register(name, parser_class, metadata=metadata,
                              set_as_default=set_as_default)

        # Create instance and store
        instance = parser_class()
        self._parsers[name] = instance

    def get_parser(self, file_path: str) -> Optional[DocumentParser]:
        """
        Get appropriate parser for a file.

        Args:
            file_path: Path to the file

        Returns:
            Parser instance or None
        """
        ext = os.path.splitext(file_path)[1].lower()

        # Check all parsers for support
        for name, parser in self._parsers.items():
            if parser.supports_file(file_path):
                return parser

        # Fall back to text parser
        text_parser = self._parsers.get('text')
        if text_parser and ext in text_parser.TEXT_EXTENSIONS:
            return text_parser

        return None

    def parse(self, file_path: str) -> str:
        """
        Parse a document using the appropriate parser.

        Args:
            file_path: Path to the document

        Returns:
            Extracted text content
        """
        parser = self.get_parser(file_path)
        if parser:
            return parser.parse(file_path)
        return ""

    def set_ocr_processor(self, ocr_processor) -> None:
        """Set OCR processor for image parser."""
        image_parser = self._parsers.get('image')
        if image_parser:
            image_parser.set_ocr_processor(ocr_processor)

    def get_supported_extensions(self) -> set:
        """Get all supported file extensions."""
        extensions = set()
        for parser in self._parsers.values():
            if hasattr(parser, 'TEXT_EXTENSIONS'):
                extensions.update(parser.TEXT_EXTENSIONS)

        for item in self._factory.get_registered_names():
            exts = self._factory.get_metadata(item, 'extensions', [])
            extensions.update(exts)

        return extensions

    def is_supported(self, file_path: str) -> bool:
        """Check if file type is supported."""
        return self.get_parser(file_path) is not None


# Global factory instance
_document_parser_factory: Optional[DocumentParserFactory] = None


def get_document_parser_factory() -> DocumentParserFactory:
    """Get the global document parser factory instance."""
    global _document_parser_factory
    if _document_parser_factory is None:
        _document_parser_factory = DocumentParserFactory()
    return _document_parser_factory


def parse_document(file_path: str) -> str:
    """Convenience function to parse a document."""
    return get_document_parser_factory().parse(file_path)
