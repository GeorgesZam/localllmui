import os
import json
import re
import shutil
import hashlib
import numpy as np
from typing import Optional, Callable, List, Tuple

import config
from utils import get_resource_path, get_writable_path, get_file_hash
from ocr import OCRProcessor

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import openpyxl
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False


class EmbeddingModel:

    def __init__(self):
        self._model = None
        self._loaded = False

    @property
    def model(self):
        return self._model

    @property
    def is_loaded(self) -> bool:
        return self._loaded

    def load(self, on_progress: Optional[Callable[[str], None]] = None) -> bool:
        if self._loaded:
            return True

        def log(msg):
            print(f"[Embedding] {msg}")
            if on_progress:
                on_progress(msg)

        try:
            from sentence_transformers import SentenceTransformer
            bundled_path = get_resource_path(config.EMBEDDING_MODEL_FOLDER)

            if os.path.exists(bundled_path):
                log(f"Loading: {bundled_path}")

                # Add timeout for model loading to prevent hanging on Windows
                import signal
                import threading

                load_result = {'success': False, 'error': None}

                def load_with_timeout():
                    try:
                        self._model = SentenceTransformer(bundled_path)
                        load_result['success'] = True
                    except Exception as e:
                        load_result['error'] = str(e)

                # Start loading in a thread
                load_thread = threading.Thread(target=load_with_timeout, daemon=True)
                load_thread.start()

                # Wait with timeout (30 seconds should be enough)
                load_thread.join(timeout=30)

                if load_thread.is_alive():
                    log("Loading timeout - model may be too large or corrupted")
                    return False

                if not load_result['success']:
                    log(f"Error loading model: {load_result['error']}")
                    return False

                self._loaded = True
                log("Embedding model loaded!")
                return True
            else:
                log(f"Model not found: {bundled_path}")
                return False
        except Exception as e:
            log(f"Error: {e}")
            import traceback
            traceback.print_exc()
            return False

    def encode(self, texts: List[str], is_query: bool = False) -> np.ndarray:
        if self._model is None:
            return np.array([])

        if is_query:
            texts = [f"Represent this sentence for searching relevant passages: {t}" for t in texts]

        batch_size = getattr(config, 'BATCH_SIZE', 512)
        return self._model.encode(
            texts,
            normalize_embeddings=True,
            show_progress_bar=False,
            batch_size=batch_size
        )

    def unload(self):
        self._model = None
        self._loaded = False
        import gc
        gc.collect()


class DocumentParser:

    SUPPORTED_EXTENSIONS = (
        '.txt', '.md', '.pdf', '.xlsx', '.xls', '.pptx', '.ppt', '.csv',
        '.docx', '.doc', '.py', '.js', '.ts', '.jsx', '.tsx', '.java',
        '.c', '.cpp', '.h', '.hpp', '.cs', '.go', '.rs', '.rb', '.php',
        '.json', '.xml', '.yaml', '.yml', '.html', '.htm', '.css',
        '.png', '.jpg', '.jpeg', '.tiff', '.bmp',
    )

    TEXT_EXTENSIONS = (
        '.txt', '.md', '.py', '.js', '.ts', '.jsx', '.tsx', '.java',
        '.c', '.cpp', '.h', '.hpp', '.cs', '.go', '.rs', '.rb', '.php',
        '.json', '.xml', '.yaml', '.yml', '.html', '.htm', '.css', '.csv',
    )

    def __init__(self, ocr_processor: OCRProcessor):
        self.ocr = ocr_processor
        self._parsers = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx, '.doc': self._parse_docx,
            '.xlsx': self._parse_excel, '.xls': self._parse_excel,
            '.pptx': self._parse_pptx, '.ppt': self._parse_pptx,
            '.png': self._parse_image, '.jpg': self._parse_image,
            '.jpeg': self._parse_image, '.tiff': self._parse_image,
            '.bmp': self._parse_image,
        }

    def parse(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        print(f"[Parser] Reading: {os.path.basename(file_path)} ({ext})")

        parser = self._parsers.get(ext, self._parse_text)
        return parser(file_path)

    def _parse_pdf(self, file_path: str) -> str:
        if not HAS_PDF:
            return ""

        try:
            text_parts = []
            scanned_pages = []

            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)

                for page_num, page in enumerate(reader.pages):
                    try:
                        page_text = (page.extract_text() or "").strip()

                        if len(page_text) > 50:
                            text_parts.append(f"=== Page {page_num + 1} ===\n{page_text}")
                        else:
                            scanned_pages.append(page_num + 1)
                    except Exception:
                        scanned_pages.append(page_num + 1)

            text = "\n\n".join(text_parts)

            if not text.strip() and scanned_pages and self.ocr.available and self.ocr.pdf_support:
                text = self.ocr.ocr_pdf(file_path, dpi=200)

            return text.strip()
        except Exception as e:
            print(f"[Parser] PDF error: {e}")
            return ""

    def _parse_docx(self, file_path: str) -> str:
        if not HAS_DOCX:
            return ""

        try:
            doc = Document(file_path)
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]

            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        text_parts.append(row_text)

            return "\n".join(text_parts)
        except Exception as e:
            print(f"[Parser] DOCX error: {e}")
            return ""

    def _parse_excel(self, file_path: str) -> str:
        if not HAS_EXCEL:
            return ""

        try:
            text_parts = []
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) if c else "" for c in row)
                    if row_text.strip():
                        text_parts.append(row_text)

            wb.close()
            return "\n".join(text_parts)
        except Exception as e:
            print(f"[Parser] Excel error: {e}")
            return ""

    def _parse_pptx(self, file_path: str) -> str:
        if not HAS_PPTX:
            return ""

        try:
            prs = Presentation(file_path)
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = [s.text for s in slide.shapes if hasattr(s, "text") and s.text]
                if slide_text:
                    text_parts.append(f"=== Slide {i} ===\n" + "\n".join(slide_text))

            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"[Parser] PPTX error: {e}")
            return ""

    def _parse_image(self, file_path: str) -> str:
        if not self.ocr.available:
            return ""
        return self.ocr.ocr_image(file_path)

    def _parse_text(self, file_path: str) -> str:
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception:
                return ""
        return ""


class RAG:

    def __init__(self):
        self.documents = []
        self._embeddings = None
        self.embedding_model = EmbeddingModel()
        self.ocr_processor = OCRProcessor()
        self.parser = DocumentParser(self.ocr_processor)
        self.chunks = {}  # Store chunks by conversation_id
        self.conversation_docs = {}  # Store document mappings by conversation_id
        self.last_sources = []
        self._index_hash = None

        self.index_file = get_writable_path("index.json")
        self.embeddings_file = get_writable_path("embeddings.npy")
        self.hash_file = get_writable_path("index_hash.txt")
        self.user_docs_folder = get_writable_path("documents")
        self.bundled_data_folder = get_resource_path(config.RAG_FOLDER)

    @property
    def embeddings(self):
        if self._embeddings is None and os.path.exists(self.embeddings_file):
            try:
                self._embeddings = np.load(self.embeddings_file)
            except Exception:
                pass
        return self._embeddings

    @embeddings.setter
    def embeddings(self, value):
        self._embeddings = value

    def initialize(self, on_progress: Optional[Callable[[str], None]] = None) -> bool:
        def log(msg):
            print(f"[RAG] {msg}")
            if on_progress:
                on_progress(msg)

        if not config.RAG_ENABLED:
            return True

        os.makedirs(self.user_docs_folder, exist_ok=True)

        ocr_status = self.ocr_processor.get_status()
        log("✓ OCR available" if ocr_status["ocr_available"] else "⚠ OCR not available")

        log("Loading embedding model...")
        if not self.embedding_model.load(on_progress):
            log("⚠ Using keyword search (no embedding)")

        if self._is_cache_valid():
            log("Loading cached index...")
            self._load_index()
        else:
            log("Building new index...")
            self._build_index(log)

        log(f"RAG ready: {len(self.documents)} chunks")
        return True

    def _compute_docs_hash(self) -> str:
        hashes = []

        for folder in [self.bundled_data_folder, self.user_docs_folder]:
            if not os.path.isdir(folder):
                continue
            for filename in sorted(os.listdir(folder)):
                ext = os.path.splitext(filename)[1].lower()
                if ext in DocumentParser.SUPPORTED_EXTENSIONS:
                    filepath = os.path.join(folder, filename)
                    hashes.append(get_file_hash(filepath))

        return hashlib.md5("|".join(hashes).encode()).hexdigest()

    def _is_cache_valid(self) -> bool:
        if not config.INDEX_CACHE_ENABLED:
            return False

        if not os.path.exists(self.index_file):
            return False

        if not os.path.exists(self.hash_file):
            return False

        try:
            with open(self.hash_file, 'r') as f:
                cached_hash = f.read().strip()
            return cached_hash == self._compute_docs_hash()
        except Exception:
            return False

    def _load_index(self):
        try:
            with open(self.index_file, "r", encoding="utf-8") as f:
                self.documents = json.load(f)

            self._embeddings = None
        except Exception as e:
            print(f"[RAG] Load error: {e}")
            self.documents = []
            self._embeddings = None

    def _save_index(self, log):
        try:
            with open(self.index_file, "w", encoding="utf-8") as f:
                json.dump(self.documents, f, ensure_ascii=False)

            if self._embeddings is not None:
                np.save(self.embeddings_file, self._embeddings)

            with open(self.hash_file, 'w') as f:
                f.write(self._compute_docs_hash())

            log(f"Index saved: {len(self.documents)} chunks")
        except Exception as e:
            log(f"Save error: {e}")

    def _split_text(self, text: str, chunk_size: int) -> List[str]:
        overlap = getattr(config, 'RAG_CHUNK_OVERLAP', 50)
        text = text.strip()

        if len(text) < 50:
            return [text] if len(text) > 20 else []

        sentences = re.split(r'(?<=[.!?])\s+', text)
        chunks, current_chunk, current_length = [], [], 0

        for sentence in sentences:
            words = sentence.split()
            if current_length + len(words) <= chunk_size:
                current_chunk.extend(words)
                current_length += len(words)
            else:
                if current_chunk:
                    chunks.append(" ".join(current_chunk))
                overlap_words = current_chunk[-overlap:] if overlap > 0 else []
                current_chunk = overlap_words + words
                current_length = len(current_chunk)

        if current_chunk and len(current_chunk) >= 15:
            chunks.append(" ".join(current_chunk))

        return chunks

    def _build_index(self, log):
        all_chunks = []

        folders = []
        if os.path.exists(self.bundled_data_folder):
            folders.append(("bundled", self.bundled_data_folder))
        if os.path.exists(self.user_docs_folder):
            folders.append(("user", self.user_docs_folder))

        for folder_type, folder in folders:
            log(f"Scanning {folder_type}: {folder}")

            if not os.path.isdir(folder):
                continue

            for filename in os.listdir(folder):
                ext = os.path.splitext(filename)[1].lower()
                if ext not in DocumentParser.SUPPORTED_EXTENSIONS:
                    continue

                file_path = os.path.join(folder, filename)

                try:
                    log(f"Processing: {filename}")
                    text = self.parser.parse(file_path)

                    if not text or len(text.strip()) < 10:
                        log(f"⚠ {filename}: empty")
                        continue

                    chunks = self._split_text(text, config.RAG_CHUNK_SIZE)

                    for i, chunk in enumerate(chunks):
                        all_chunks.append({
                            "source": filename,
                            "chunk_id": i,
                            "content": chunk
                        })

                    log(f"✓ {filename}: {len(chunks)} chunks")
                except Exception as e:
                    log(f"✗ {filename}: {e}")

        if not all_chunks:
            log("No documents found")
            self.documents = []
            self._embeddings = None
            return

        if self.embedding_model.is_loaded:
            log(f"Encoding {len(all_chunks)} chunks...")
            texts = [c["content"] for c in all_chunks]
            self._embeddings = self.embedding_model.encode(texts, is_query=False)
        else:
            self._embeddings = None

        self.documents = all_chunks
        self._save_index(log)

    def search(self, query: str, top_k: int = None, max_context_chars: int = 1200, allowed_sources: List[str] = None) -> Tuple[str, List[dict]]:
        """
        Search for relevant documents with context size limit.

        Args:
            query: Search query
            top_k: Number of results to return
            max_context_chars: Maximum characters for context (to avoid context window overflow)
            allowed_sources: Optional list of document filenames to filter results (for conversation isolation)

        Returns:
            Tuple of (context string, list of sources)
        """
        # If allowed_sources is provided, create a set for faster lookup
        # None means search all documents, [] means search no documents
        allowed_set = set(allowed_sources) if allowed_sources is not None else None
        if not self.documents:
            self.last_sources = []
            return "", []

        top_k = top_k or config.RAG_TOP_K
        min_score = getattr(config, 'RAG_MIN_SCORE', 0.3)
        results = []

        # Filter documents by allowed sources if provided
        searchable_docs = self.documents
        if allowed_set is not None:
            searchable_docs = [d for d in self.documents if d["source"] in allowed_set]

        if not searchable_docs:
            self.last_sources = []
            return "", []

        if self.embeddings is not None and self.embedding_model.is_loaded:
            # Create a mapping from original document indices to filtered ones
            original_to_filtered = {i: doc for i, doc in enumerate(self.documents)}
            filtered_indices = [i for i, doc in enumerate(self.documents) if doc["source"] in allowed_set] if allowed_set is not None else list(range(len(self.documents)))

            # Get embeddings for filtered documents only
            filtered_embeddings = self.embeddings[filtered_indices] if allowed_set is not None else self.embeddings
            filtered_docs = [self.documents[i] for i in filtered_indices] if allowed_set is not None else self.documents

            query_emb = self.embedding_model.encode([query], is_query=True)[0]
            similarities = np.dot(filtered_embeddings, query_emb)
            top_indices = np.argsort(similarities)[-top_k * 2:][::-1]

            for idx in top_indices:
                score = float(similarities[idx])
                if score >= min_score:
                    results.append((filtered_docs[idx], score))
            results = results[:top_k]
        else:
            query_words = set(query.lower().split())
            scored = []
            for doc in searchable_docs:
                content_lower = doc["content"].lower()
                content_words = set(content_lower.split())
                matches = query_words & content_words
                score = len(matches) + sum(0.5 for w in query_words if w in content_lower)
                if score > 0:
                    scored.append((doc, score))
            scored.sort(key=lambda x: x[1], reverse=True)
            results = scored[:top_k]

        if not results:
            self.last_sources = []
            return "", []

        self.last_sources = []
        context_parts = []

        for i, (doc, score) in enumerate(results, 1):
            content = doc["content"]

            # Calculate header length
            header = f"[Document {i} - {doc['source']}]\n"
            header_len = len(header)

            # Reserve space for header and "..."
            max_content_len = max_context_chars // len(results) - header_len - 10

            # Truncate if needed
            if len(content) > max_content_len:
                content = content[:max_content_len] + "..."

            self.last_sources.append({
                "index": i,
                "source": doc["source"],
                "chunk_id": doc["chunk_id"],
                "score": score,
                "preview": doc["content"][:200] + "..." if len(doc["content"]) > 200 else doc["content"]
            })

            context_parts.append(f"{header}{content}")

        return "\n\n".join(context_parts), self.last_sources

    def format_sources_for_display(self) -> str:
        if not self.last_sources:
            return ""

        lines = ["", "📚 Sources:"]
        for src in self.last_sources:
            lines.append(f"  [{src['index']}] {src['source']} (score: {src['score']:.2f})")
            preview = src['preview'][:80].replace('\n', ' ')
            lines.append(f"      \"{preview}...\"")

        return "\n".join(lines)

    def add_documents(self, file_paths: list, on_progress: Optional[Callable[[str], None]] = None) -> bool:
        def log(msg):
            print(f"[RAG] {msg}")
            if on_progress:
                on_progress(msg)

        try:
            os.makedirs(self.user_docs_folder, exist_ok=True)
            added = 0

            for file_path in file_paths:
                filename = os.path.basename(file_path)
                ext = os.path.splitext(filename)[1].lower()

                if ext not in DocumentParser.SUPPORTED_EXTENSIONS:
                    log(f"⚠ {filename}: unsupported")
                    continue

                dest = os.path.join(self.user_docs_folder, filename)
                shutil.copy2(file_path, dest)
                log(f"✓ Copied: {filename}")
                added += 1

            if added == 0:
                return False

            log("Rebuilding index...")
            self._build_index(log)

            return len(self.documents) > 0
        except Exception as e:
            log(f"Error: {e}")
            return False

    def clear_cache(self):
        for f in [self.index_file, self.embeddings_file, self.hash_file]:
            if os.path.exists(f):
                os.remove(f)
        self.documents = []
        self._embeddings = None

    def add_conversation_mapping(self, conversation_id: str, document_ids: List[str]):
        """Add conversation to document mapping."""
        if not hasattr(self, 'conversation_docs'):
            self.conversation_docs = {}
        self.conversation_docs[conversation_id] = document_ids

    def get_conversation_documents(self, conversation_id: str) -> List[str]:
        """Get documents for a specific conversation."""
        if not hasattr(self, 'conversation_docs'):
            return []
        return self.conversation_docs.get(conversation_id, [])

    def add_document_to_conversation(self, conversation_id: str, document_id: str):
        """Add a document to a conversation."""
        if not hasattr(self, 'conversation_docs'):
            self.conversation_docs = {}

        if conversation_id not in self.conversation_docs:
            self.conversation_docs[conversation_id] = []

        if document_id not in self.conversation_docs[conversation_id]:
            self.conversation_docs[conversation_id].append(document_id)

    def remove_conversation(self, conversation_id: str):
        """Remove a conversation and its document mappings."""
        if hasattr(self, 'conversation_docs'):
            self.conversation_docs.pop(conversation_id, None)
