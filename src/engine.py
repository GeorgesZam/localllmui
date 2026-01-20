"""
LLM Engine with RAG (embeddings) and conversation history.
FIXED: Loads bundled embedding model, no downloads at runtime.
"""

import os
import sys
import json
import numpy as np
from typing import Iterator, Optional, Callable, List, Tuple
import config

# Document parsers
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import openpyxl
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def get_resource_path(relative_path: str) -> str:
    """Gets path compatible with PyInstaller (read-only bundled files)."""
    if hasattr(sys, '_MEIPASS'):
        return os.path.join(sys._MEIPASS, relative_path)
    return relative_path


def get_writable_path(filename: str) -> str:
    """Gets writable path for user data - works with PyInstaller."""
    if hasattr(sys, '_MEIPASS'):
        app_data = os.path.join(os.path.expanduser("~"), ".localchat")
    else:
        app_data = os.path.join(os.path.dirname(__file__), "..", ".localchat")
    
    os.makedirs(app_data, exist_ok=True)
    return os.path.join(app_data, filename)


class EmbeddingModel:
    """Embedding model using sentence-transformers - loads from bundled folder."""
    
    def __init__(self):
        self.model = None
    
    def load(self, on_progress: Optional[Callable[[str], None]] = None):
        """Loads the embedding model from bundled files (no download)."""
        def log(msg):
            print(f"[Embedding] {msg}")
            if on_progress:
                on_progress(msg)
        
        try:
            from sentence_transformers import SentenceTransformer
            
            # Load from bundled folder - NO DOWNLOAD
            bundled_model_path = get_resource_path(config.EMBEDDING_MODEL_FOLDER)
            
            if os.path.exists(bundled_model_path):
                log(f"Loading bundled embedding model: {bundled_model_path}")
                self.model = SentenceTransformer(bundled_model_path)
                log("Embedding model loaded!")
                return True
            else:
                log(f"Bundled model not found at: {bundled_model_path}")
                log("Falling back to keyword search")
                return False
                
        except Exception as e:
            log(f"Error loading embedding model: {e}")
            return False
    
    def encode(self, texts: List[str], is_query: bool = False) -> np.ndarray:
        """Encodes texts to vectors."""
        if self.model is None:
            return np.array([])
        
        # BGE models work better with query prefix
        if is_query:
            texts = [f"Represent this sentence for searching relevant passages: {t}" for t in texts]
        
        return self.model.encode(texts, normalize_embeddings=True, show_progress_bar=False)


class RAG:
    """RAG with semantic search and source tracking."""
    
    SUPPORTED_EXTENSIONS = (
        '.txt', '.md', '.pdf', '.xlsx', '.xls', '.pptx', '.ppt', '.csv',
        '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.c', '.cpp', '.h', '.hpp',
        '.cs', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.scala', '.r',
        '.json', '.xml', '.yaml', '.yml', '.ini', '.cfg', '.conf', '.toml',
        '.sh', '.bash', '.zsh', '.ps1', '.bat', '.sql',
        '.html', '.htm', '.css', '.scss', '.sass', '.less',
    )
    
    def __init__(self):
        self.documents = []
        self.embeddings = None
        self.embedding_model = EmbeddingModel()
        self.last_sources = []
        
        self.index_file = get_writable_path("index.json")
        self.embeddings_file = get_writable_path("embeddings.npy")
        self.user_docs_folder = get_writable_path("documents")
        self.bundled_data_folder = get_resource_path(config.RAG_FOLDER)
    
    def initialize(self, on_progress: Optional[Callable[[str], None]] = None) -> bool:
        """Initializes RAG system."""
        def log(msg):
            print(f"[RAG] {msg}")
            if on_progress:
                on_progress(msg)
        
        if not config.RAG_ENABLED:
            return True
        
        os.makedirs(self.user_docs_folder, exist_ok=True)
        log(f"User docs folder: {self.user_docs_folder}")
        
        log("Loading embedding model...")
        if not self.embedding_model.load(on_progress):
            log("Embedding model failed, using keyword search")
        
        if os.path.exists(self.index_file) and os.path.exists(self.embeddings_file):
            log("Loading existing index...")
            self._load_index()
            if self.documents:
                log(f"Loaded {len(self.documents)} chunks from index")
            else:
                log("Index empty, rebuilding...")
                self._build_index(log)
        else:
            log("No index found, building...")
            self._build_index(log)
        
        log(f"RAG ready: {len(self.documents)} chunks indexed")
        return True
    
    def _load_index(self):
        """Loads pre-built index."""
        try:
            with open(self.index_file, "r", encoding="utf-8") as f:
                self.documents = json.load(f)
            
            if os.path.exists(self.embeddings_file) and self.embedding_model.model is not None:
                self.embeddings = np.load(self.embeddings_file)
            else:
                self.embeddings = None
                
            print(f"[RAG] Index loaded: {len(self.documents)} docs")
        except Exception as e:
            print(f"[RAG] Error loading index: {e}")
            self.documents = []
            self.embeddings = None
    
    def _save_index(self, log: Callable[[str], None]):
        """Saves index to disk."""
        try:
            with open(self.index_file, "w", encoding="utf-8") as f:
                json.dump(self.documents, f, ensure_ascii=False, indent=2)
            if self.embeddings is not None:
                np.save(self.embeddings_file, self.embeddings)
            log(f"Index saved: {len(self.documents)} chunks")
        except Exception as e:
            log(f"Failed to save index: {e}")
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        if not HAS_PDF:
            return ""
        try:
            text = ""
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text.strip()
        except Exception as e:
            print(f"[RAG] PDF error {file_path}: {e}")
            return ""
    
    def _extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from Excel."""
        if not HAS_EXCEL:
            return ""
        try:
            text = ""
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n=== Sheet: {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            wb.close()
            return text.strip()
        except Exception as e:
            print(f"[RAG] Excel error {file_path}: {e}")
            return ""
    
    def _extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint."""
        if not HAS_PPTX:
            return ""
        try:
            text = ""
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n=== Slide {i} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"[RAG] PPTX error {file_path}: {e}")
            return ""
    
    def _read_document(self, file_path: str) -> str:
        """Read document based on extension."""
        ext = os.path.splitext(file_path)[1].lower()
        
        print(f"[RAG] Reading: {file_path} (ext: {ext})")
        
        if ext == '.pdf':
            return self._extract_text_from_pdf(file_path)
        elif ext in ('.xlsx', '.xls'):
            return self._extract_text_from_excel(file_path)
        elif ext in ('.pptx', '.ppt'):
            return self._extract_text_from_pptx(file_path)
        else:
            encodings = ['utf-8', 'latin-1', 'cp1252']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    print(f"[RAG] Read {len(content)} chars from {file_path}")
                    return content
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    print(f"[RAG] Error reading {file_path}: {e}")
                    return ""
            return ""
    
    def _split_text(self, text: str, chunk_size: int, overlap: int = None) -> List[str]:
        """Splits text into overlapping chunks."""
        if overlap is None:
            overlap = getattr(config, 'RAG_CHUNK_OVERLAP', 100)
        
        text = text.strip()
        if len(text) < 50:
            return [text] if len(text) > 20 else []
        
        import re
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence_words = sentence.split()
            sentence_length = len(sentence_words)
            
            if current_length + sentence_length <= chunk_size:
                current_chunk.extend(sentence_words)
                current_length += sentence_length
            else:
                if current_chunk:
                    chunks.append(" ".join(current_chunk))
                
                if overlap > 0 and current_chunk:
                    overlap_words = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                    current_chunk = overlap_words + sentence_words
                    current_length = len(current_chunk)
                else:
                    current_chunk = sentence_words
                    current_length = sentence_length
        
        if current_chunk and len(current_chunk) >= 20:
            chunks.append(" ".join(current_chunk))
        
        if not chunks and len(text.split()) >= 20:
            words = text.split()
            step = max(1, chunk_size - overlap)
            for i in range(0, len(words), step):
                chunk_words = words[i:i + chunk_size]
                if len(chunk_words) >= 20:
                    chunks.append(" ".join(chunk_words))
        
        return chunks
    
    def _build_index(self, log: Callable[[str], None]):
        """Builds index from all document folders."""
        all_chunks = []
        
        folders = []
        if os.path.exists(self.bundled_data_folder):
            folders.append(("bundled", self.bundled_data_folder))
        if os.path.exists(self.user_docs_folder):
            folders.append(("user", self.user_docs_folder))
        
        for folder_type, folder in folders:
            log(f"Scanning {folder_type} folder: {folder}")
            
            if not os.path.isdir(folder):
                continue
                
            for filename in os.listdir(folder):
                ext = os.path.splitext(filename)[1].lower()
                
                if ext not in self.SUPPORTED_EXTENSIONS:
                    log(f"⏭ Skip {filename} (unsupported: {ext})")
                    continue
                
                file_path = os.path.join(folder, filename)
                
                try:
                    text = self._read_document(file_path)
                    
                    if not text or len(text.strip()) < 10:
                        log(f"⚠ {filename}: empty or too short")
                        continue
                    
                    chunks = self._split_text(text, config.RAG_CHUNK_SIZE)
                    
                    for i, chunk in enumerate(chunks):
                        all_chunks.append({
                            "source": filename,
                            "chunk_id": i,
                            "content": chunk
                        })
                    
                    log(f"✓ {filename}: {len(chunks)} chunks ({len(text)} chars)")
                    
                except Exception as e:
                    log(f"✗ {filename}: {e}")
        
        if not all_chunks:
            log("No documents found to index")
            self.documents = []
            self.embeddings = None
            return
        
        # Only encode if embedding model is available
        if self.embedding_model.model is not None:
            log(f"Encoding {len(all_chunks)} chunks...")
            texts = [c["content"] for c in all_chunks]
            self.embeddings = self.embedding_model.encode(texts, is_query=False)
        else:
            log(f"Indexing {len(all_chunks)} chunks (keyword mode)...")
            self.embeddings = None
        
        self.documents = all_chunks
        self._save_index(log)
    
    def search(self, query: str, top_k: int = None) -> Tuple[str, List[dict]]:
        """Searches documents using semantic or keyword search."""
        if not self.documents:
            print("[RAG] No documents indexed!")
            self.last_sources = []
            return "", []
        
        top_k = top_k or config.RAG_TOP_K
        min_score = getattr(config, 'RAG_MIN_SCORE', 0.25)
        
        print(f"[RAG] Searching '{query}' in {len(self.documents)} chunks")
        
        results = []
        
        if self.embeddings is not None and self.embedding_model.model is not None:
            # Semantic search
            query_embedding = self.embedding_model.encode([query], is_query=True)[0]
            similarities = np.dot(self.embeddings, query_embedding)
            
            top_indices = np.argsort(similarities)[-top_k * 2:][::-1]
            
            for idx in top_indices:
                score = float(similarities[idx])
                if score >= min_score:
                    results.append((self.documents[idx], score))
                    print(f"[RAG] ✓ {self.documents[idx]['source']} chunk {self.documents[idx]['chunk_id']} (score: {score:.3f})")
            
            results = results[:top_k]
        else:
            # Keyword fallback
            query_words = set(query.lower().split())
            scored = []
            
            for doc in self.documents:
                content_lower = doc["content"].lower()
                content_words = set(content_lower.split())
                
                matches = query_words & content_words
                score = len(matches)
                
                for word in query_words:
                    if word in content_lower:
                        score += 0.5
                
                if score > 0:
                    scored.append((doc, score))
            
            scored.sort(key=lambda x: x[1], reverse=True)
            results = scored[:top_k]
            
            for doc, score in results:
                print(f"[RAG] ✓ {doc['source']} chunk {doc['chunk_id']} (keyword score: {score:.1f})")
        
        if not results:
            print("[RAG] No matches found")
            self.last_sources = []
            return "", []
        
        # Store sources for display
        self.last_sources = []
        context_parts = []
        
        for i, (doc, score) in enumerate(results, 1):
            self.last_sources.append({
                "index": i,
                "source": doc["source"],
                "chunk_id": doc["chunk_id"],
                "score": score,
                "preview": doc["content"][:200] + "..." if len(doc["content"]) > 200 else doc["content"]
            })
            context_parts.append(f"[Document {i} - {doc['source']}]\n{doc['content']}")
        
        context = "\n\n".join(context_parts)
        print(f"[RAG] Returning {len(results)} chunks, {len(context)} chars")
        
        return context, self.last_sources
    
    def get_last_sources(self) -> List[dict]:
        """Returns the sources used in the last search."""
        return self.last_sources
    
    def format_sources_for_display(self) -> str:
        """Formats sources for UI display."""
        if not self.last_sources:
            return ""
        
        lines = ["", "📚 Sources used:"]
        for src in self.last_sources:
            lines.append(f"  [{src['index']}] {src['source']} (chunk {src['chunk_id']}, score: {src['score']:.2f})")
            preview = src['preview'][:150].replace('\n', ' ')
            lines.append(f"      \"{preview}...\"")
        
        return "\n".join(lines)
    
    def add_documents(self, file_paths: list, on_progress: Optional[Callable[[str], None]] = None) -> bool:
        """Adds new documents and rebuilds index."""
        def log(msg):
            print(f"[RAG] {msg}")
            if on_progress:
                on_progress(msg)
        
        try:
            os.makedirs(self.user_docs_folder, exist_ok=True)
            
            import shutil
            added_count = 0
            
            for file_path in file_paths:
                filename = os.path.basename(file_path)
                ext = os.path.splitext(filename)[1].lower()
                
                if ext not in self.SUPPORTED_EXTENSIONS:
                    log(f"⚠ {filename}: unsupported ({ext})")
                    continue
                
                dest = os.path.join(self.user_docs_folder, filename)
                shutil.copy2(file_path, dest)
                log(f"✓ Copied: {filename}")
                added_count += 1
            
            if added_count == 0:
                log("No supported files to add")
                return False
            
            log("Rebuilding index...")
            self._build_index(log)
            
            return len(self.documents) > 0
            
        except Exception as e:
            log(f"Error: {e}")
            import traceback
            traceback.print_exc()
            return False


class LLMEngine:
    """LLM Engine with history and RAG."""
    
    def __init__(self):
        self.llm = None
        self.history = []
        self.rag = RAG()
        self.is_ready = False
        self.error = None
    
    def load(self, on_progress: Optional[Callable[[str], None]] = None) -> bool:
        """Loads the LLM model and RAG."""
        def log(msg):
            print(f"[Engine] {msg}")
            if on_progress:
                on_progress(msg)
        
        try:
            self.rag.initialize(log)
            
            log("Importing llama_cpp...")
            from llama_cpp import Llama
            
            model_path = get_resource_path(config.MODEL_FILE)
            log(f"Model: {model_path}")
            
            if not os.path.exists(model_path):
                raise FileNotFoundError(f"Model not found: {model_path}")
            
            log("Loading LLM...")
            self.llm = Llama(
                model_path=model_path,
                n_ctx=config.CONTEXT_SIZE,
                n_threads=config.THREADS,
                n_gpu_layers=-1,
                verbose=True
            )
            
            self.is_ready = True
            log("Ready!")
            return True
            
        except Exception as e:
            self.error = str(e)
            log(f"Error: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def generate(self, message: str) -> Iterator[str]:
        """Generates response with RAG and history."""
        if not self.is_ready:
            yield "Error: Model not ready"
            return
        
        # Get RAG context
        rag_context = ""
        sources = []
        if config.RAG_ENABLED and self.rag.documents:
            print(f"[Engine] RAG has {len(self.rag.documents)} documents")
            rag_context, sources = self.rag.search(message)
            print(f"[Engine] RAG context: {len(rag_context)} chars, {len(sources)} sources")
        
        # Build prompt
        if rag_context:
            system_with_context = f"""{config.SYSTEM_PROMPT}

=== CONTEXT DOCUMENTS ===
{rag_context}
=== END CONTEXT ===

Based on the context above, answer the following question. If the answer is not in the context, say so."""
        else:
            system_with_context = config.SYSTEM_PROMPT
        
        # Build conversation (Qwen format)
        prompt = f"<|im_start|>system\n{system_with_context}<|im_end|>\n"
        
        for h in self.history[-3:]:
            prompt += f"<|im_start|>user\n{h['user']}<|im_end|>\n"
            prompt += f"<|im_start|>assistant\n{h['assistant']}<|im_end|>\n"
        
        prompt += f"<|im_start|>user\n{message}<|im_end|>\n"
        prompt += "<|im_start|>assistant\n"
        
        print(f"[Engine] Prompt length: {len(prompt)} chars")
        
        # Generate
        full_response = ""
        try:
            for chunk in self.llm(
                prompt,
                max_tokens=config.MAX_TOKENS,
                stop=config.STOP_TOKENS,
                temperature=getattr(config, 'TEMPERATURE', 0.2),
                top_p=getattr(config, 'TOP_P', 0.9),
                repeat_penalty=getattr(config, 'REPEAT_PENALTY', 1.1),
                stream=True
            ):
                token = chunk["choices"][0]["text"]
                full_response += token
                yield token
        except Exception as e:
            print(f"[Engine] Generation error: {e}")
            yield f"\n[Error: {e}]"
        
        # Append sources if enabled
        if getattr(config, 'RAG_SHOW_SOURCES', True) and sources:
            sources_text = self.rag.format_sources_for_display()
            yield sources_text
            full_response += sources_text
        
        # Save to history
        if full_response.strip():
            clean_response = full_response.split("📚 Sources used:")[0].strip()
            self.history.append({
                "user": message,
                "assistant": clean_response
            })
    
    def clear_history(self):
        """Clears conversation history."""
        self.history = []
        print("[Engine] History cleared")
