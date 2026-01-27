#!/usr/bin/env python3
"""
Local Chat v4 - Fixed embedding model loading
"""

import sys
import os
import json
import re
import shutil
import multiprocessing
from datetime import datetime
from typing import Optional, List, Dict, Callable, Iterator, Tuple
from dataclasses import dataclass, asdict
from functools import lru_cache

if sys.platform == 'darwin':
    os.environ['TK_SILENCE_DEPRECATION'] = '1'
if hasattr(sys, '_MEIPASS'):
    multiprocessing.freeze_support()

# Only import what we absolutely need at startup
import numpy as np
import customtkinter as ctk
from tkinter import filedialog, messagebox
import threading

# ============================================================================
# DEBUG
# ============================================================================

DEBUG = True

def debug_log(component: str, msg: str):
    if DEBUG:
        timestamp = datetime.now().strftime("%H:%M:%S.%f")[:-3]
        print(f"[{timestamp}] [{component}] {msg}")

# ============================================================================
# CONFIG
# ============================================================================

class Config:
    APP_NAME = "Local Chat v4"
    WINDOW_SIZE = "1100x700"
    
    MODEL_FILE = "models/model.gguf"
    EMBEDDING_MODEL_DIR = "embedding_model"  # Directory name for bundled model
    EMBEDDING_MODEL_NAME = "BAAI/bge-small-en-v1.5"  # Fallback HuggingFace model
    
    _CPU_COUNT = multiprocessing.cpu_count()
    CONTEXT_SIZE = 4096
    MAX_TOKENS = 512
    THREADS = max(4, _CPU_COUNT - 2)
    GPU_LAYERS = -1
    
    SYSTEM_PROMPT = """You are a helpful assistant. Answer questions based ONLY on the provided context documents.
If the answer is not found in the context, say "I don't have this information in the provided documents."
Be concise and specific. Quote relevant parts when possible.
Answer in the same language as the user."""
    
    STOP_TOKENS = ["<|im_end|>", "<end_of_turn>", "<|endoftext|>"]
    
    RAG_ENABLED = True
    RAG_CHUNK_SIZE = 400
    RAG_CHUNK_OVERLAP = 50
    RAG_TOP_K = 5
    RAG_MIN_SCORE = 0.25
    RAG_SHOW_SOURCES = True
    
    TEMPERATURE = 0.2
    TOP_P = 0.9
    REPEAT_PENALTY = 1.1

config = Config()

# ============================================================================
# UTILITIES
# ============================================================================

def get_resource_path(relative_path: str) -> str:
    """Get absolute path to resource, works for dev and for PyInstaller."""
    if hasattr(sys, '_MEIPASS'):
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
        debug_log("PATH", f"PyInstaller mode: _MEIPASS={base_path}")
    else:
        # Running in development
        base_path = os.path.dirname(os.path.abspath(__file__))
        # Go up one level from src/ to project root
        base_path = os.path.dirname(base_path)
        debug_log("PATH", f"Dev mode: base={base_path}")
    
    full_path = os.path.join(base_path, relative_path)
    debug_log("PATH", f"Resource path: {relative_path} -> {full_path}")
    return full_path

def list_dir_recursive(path: str, prefix: str = "") -> List[str]:
    """List directory contents recursively for debugging."""
    results = []
    try:
        if os.path.isdir(path):
            for item in os.listdir(path):
                item_path = os.path.join(path, item)
                if os.path.isdir(item_path):
                    results.append(f"{prefix}[DIR] {item}/")
                    results.extend(list_dir_recursive(item_path, prefix + "  "))
                else:
                    size = os.path.getsize(item_path)
                    results.append(f"{prefix}[FILE] {item} ({size:,} bytes)")
    except Exception as e:
        results.append(f"{prefix}[ERROR] {e}")
    return results

_APP_DATA_DIR = None

def get_writable_path(filename: str) -> str:
    """Get path in writable app data directory."""
    global _APP_DATA_DIR
    if _APP_DATA_DIR is None:
        if hasattr(sys, '_MEIPASS'):
            _APP_DATA_DIR = os.path.join(os.path.expanduser("~"), ".localchat")
        else:
            _APP_DATA_DIR = os.path.join(os.path.dirname(__file__), "..", ".localchat")
        os.makedirs(_APP_DATA_DIR, exist_ok=True)
        debug_log("UTILS", f"App data: {_APP_DATA_DIR}")
    return os.path.join(_APP_DATA_DIR, filename)

# ============================================================================
# DOCUMENT PARSER - LAZY IMPORTS
# ============================================================================

class DocumentParser:
    SUPPORTED_EXTENSIONS = (
        '.txt', '.md', '.pdf', '.xlsx', '.xls', '.pptx', '.ppt', '.csv',
        '.docx', '.doc', '.py', '.js', '.json', '.xml', '.yaml', '.yml',
        '.html', '.css',
    )
    
    def __init__(self):
        self._pdf_reader = None
        self._docx_reader = None
        self._xlsx_reader = None
        self._pptx_reader = None
    
    def parse(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        debug_log("PARSER", f"Parsing {os.path.basename(file_path)}")
        
        if ext == '.pdf':
            return self._parse_pdf(file_path)
        elif ext in ('.docx', '.doc'):
            return self._parse_docx(file_path)
        elif ext in ('.xlsx', '.xls', '.csv'):
            return self._parse_excel(file_path)
        elif ext in ('.pptx', '.ppt'):
            return self._parse_pptx(file_path)
        else:
            return self._parse_text(file_path)
    
    def _parse_pdf(self, file_path: str) -> str:
        try:
            import PyPDF2
            text_parts = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    text = (page.extract_text() or "").strip()
                    if text:
                        text_parts.append(f"=== Page {i+1} ===\n{text}")
            return "\n\n".join(text_parts)
        except Exception as e:
            debug_log("PARSER", f"PDF error: {e}")
            return ""
    
    def _parse_docx(self, file_path: str) -> str:
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            debug_log("PARSER", f"DOCX error: {e}")
            return ""
    
    def _parse_excel(self, file_path: str) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            parts = []
            for sheet in wb.sheetnames:
                parts.append(f"=== {sheet} ===")
                for row in wb[sheet].iter_rows(values_only=True):
                    row_text = " | ".join(str(c) if c else "" for c in row)
                    if row_text.strip():
                        parts.append(row_text)
            wb.close()
            return "\n".join(parts)
        except Exception as e:
            debug_log("PARSER", f"Excel error: {e}")
            return ""
    
    def _parse_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text]
                if texts:
                    parts.append(f"=== Slide {i} ===\n" + "\n".join(texts))
            return "\n\n".join(parts)
        except Exception as e:
            debug_log("PARSER", f"PPTX error: {e}")
            return ""
    
    def _parse_text(self, file_path: str) -> str:
        for enc in ['utf-8', 'latin-1', 'cp1252']:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    return f.read()
            except:
                continue
        return ""

# ============================================================================
# EMBEDDING MODEL - FIXED LOADING
# ============================================================================

class EmbeddingModel:
    """Embedding model with completely lazy loading and better path detection."""
    
    def __init__(self):
        self._model = None
        self._loaded = False
        self._load_attempted = False
    
    @property
    def is_loaded(self) -> bool:
        return self._loaded
    
    def _is_valid_model_dir(self, path: str) -> bool:
        """Check if directory contains a valid sentence-transformers model."""
        if not os.path.exists(path) or not os.path.isdir(path):
            return False
        
        files = os.listdir(path)
        
        # Check for model weights
        has_weights = any(
            f.endswith(('.bin', '.safetensors')) or f == 'pytorch_model.bin' or f == 'model.safetensors'
            for f in files
        )
        
        # Check for config
        has_config = 'config.json' in files
        
        # Check for sentence-transformers specific files
        has_st_config = any(
            f in files for f in ['config_sentence_transformers.json', 'modules.json', 'sentence_bert_config.json']
        )
        
        debug_log("EMBEDDING", f"  Path: {path}")
        debug_log("EMBEDDING", f"  Files: {files[:10]}{'...' if len(files) > 10 else ''}")
        debug_log("EMBEDDING", f"  has_weights={has_weights}, has_config={has_config}, has_st_config={has_st_config}")
        
        return has_weights or (has_config and has_st_config)
    
    def _find_model_path(self, log: Callable[[str], None]) -> Optional[str]:
        """Find the embedding model in various locations."""
        
        # List of paths to check, in order of priority
        paths_to_check = []
        
        # 1. PyInstaller bundled path (highest priority)
        if hasattr(sys, '_MEIPASS'):
            paths_to_check.append(os.path.join(sys._MEIPASS, config.EMBEDDING_MODEL_DIR))
            paths_to_check.append(os.path.join(sys._MEIPASS, "models", config.EMBEDDING_MODEL_DIR))
        
        # 2. Relative to script (for development)
        script_dir = os.path.dirname(os.path.abspath(__file__))
        project_root = os.path.dirname(script_dir)  # Go up from src/
        
        paths_to_check.extend([
            os.path.join(project_root, config.EMBEDDING_MODEL_DIR),
            os.path.join(project_root, "models", config.EMBEDDING_MODEL_DIR),
            os.path.join(script_dir, config.EMBEDDING_MODEL_DIR),
            os.path.join(script_dir, "..", config.EMBEDDING_MODEL_DIR),
        ])
        
        # 3. User cache (fallback for downloaded models)
        paths_to_check.append(get_writable_path("embedding_model_cache"))
        
        # 4. Current working directory
        paths_to_check.extend([
            os.path.join(os.getcwd(), config.EMBEDDING_MODEL_DIR),
            os.path.join(os.getcwd(), "models", config.EMBEDDING_MODEL_DIR),
        ])
        
        log("Searching for embedding model...")
        
        for path in paths_to_check:
            path = os.path.abspath(path)
            log(f"  Checking: {path}")
            
            if self._is_valid_model_dir(path):
                log(f"  ✓ Found valid model at: {path}")
                return path
            elif os.path.exists(path):
                log(f"  ✗ Path exists but not valid model")
            else:
                log(f"  ✗ Path does not exist")
        
        return None
    
    def load(self, on_progress: Optional[Callable[[str], None]] = None) -> bool:
        if self._loaded:
            return True
        
        if self._load_attempted:
            return False
        
        self._load_attempted = True
        
        def log(msg):
            debug_log("EMBEDDING", msg)
            if on_progress:
                on_progress(msg)
        
        try:
            # Import inside function to avoid global import issues
            log("Importing sentence_transformers...")
            
            try:
                from sentence_transformers import SentenceTransformer
            except ImportError as e:
                log(f"✗ sentence_transformers not available: {e}")
                return False
            
            # Debug: Show _MEIPASS contents if in PyInstaller
            if hasattr(sys, '_MEIPASS'):
                log(f"PyInstaller _MEIPASS: {sys._MEIPASS}")
                log("Contents of _MEIPASS:")
                for line in list_dir_recursive(sys._MEIPASS)[:30]:
                    log(f"  {line}")
            
            # Find the model
            model_path = self._find_model_path(log)
            
            if model_path:
                log(f"Loading model from: {model_path}")
                try:
                    self._model = SentenceTransformer(model_path)
                    self._loaded = True
                    log("✓ Model loaded successfully!")
                    
                    # Test the model
                    test_emb = self._model.encode(["test"], normalize_embeddings=True)
                    log(f"✓ Model test passed! Embedding dim: {test_emb.shape[1]}")
                    
                    return True
                except Exception as e:
                    log(f"✗ Failed to load from path: {e}")
                    import traceback
                    traceback.print_exc()
            
            # Fallback: Download from HuggingFace
            log(f"⬇ Downloading {config.EMBEDDING_MODEL_NAME}...")
            log("(~90MB, first time only)")
            
            try:
                self._model = SentenceTransformer(config.EMBEDDING_MODEL_NAME)
                self._loaded = True
                
                # Save to cache for future use
                cache_path = get_writable_path("embedding_model_cache")
                try:
                    os.makedirs(cache_path, exist_ok=True)
                    self._model.save(cache_path)
                    log(f"✓ Saved to cache: {cache_path}")
                except Exception as e:
                    log(f"⚠ Could not cache: {e}")
                
                log("✓ Embedding model ready!")
                return True
                
            except Exception as e:
                log(f"✗ Failed to download model: {e}")
                return False
            
        except Exception as e:
            log(f"✗ Error loading embedding model: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def encode(self, texts: List[str], is_query: bool = False) -> Optional[np.ndarray]:
        if not self._loaded or self._model is None:
            return None
        
        try:
            # BGE models work better with query prefix
            if is_query:
                texts = [f"Represent this sentence for searching relevant passages: {t}" for t in texts]
            
            embeddings = self._model.encode(
                texts, 
                normalize_embeddings=True, 
                show_progress_bar=False,
                convert_to_numpy=True
            )
            return embeddings
        except Exception as e:
            debug_log("EMBEDDING", f"Encode error: {e}")
            return None

# ============================================================================
# RAG SYSTEM
# ============================================================================

class RAG:
    def __init__(self):
        self.documents = []
        self._embeddings = None
        self.embedding_model = EmbeddingModel()
        self.parser = DocumentParser()
        self.last_sources = []
        self._current_conv_id = None
        self.user_docs_folder = get_writable_path("documents")
    
    def _get_conv_index_file(self, conv_id: str) -> str:
        return get_writable_path(f"index_{conv_id}.json")
    
    def _get_conv_embeddings_file(self, conv_id: str) -> str:
        return get_writable_path(f"embeddings_{conv_id}.npy")
    
    def _get_conv_docs_folder(self, conv_id: str) -> str:
        folder = os.path.join(self.user_docs_folder, conv_id)
        os.makedirs(folder, exist_ok=True)
        return folder
    
    def initialize(self, on_progress: Optional[Callable[[str], None]] = None) -> bool:
        def log(msg):
            debug_log("RAG", msg)
            if on_progress:
                on_progress(msg)
        
        os.makedirs(self.user_docs_folder, exist_ok=True)
        
        log("Loading embedding model...")
        if self.embedding_model.load(on_progress):
            log("✓ Semantic search enabled")
        else:
            log("⚠ Will use keyword search (embedding model not available)")
        
        return True
    
    def load_conversation_documents(self, conv_id: str):
        debug_log("RAG", f"Loading docs for: {conv_id}")
        self._current_conv_id = conv_id
        
        index_file = self._get_conv_index_file(conv_id)
        embeddings_file = self._get_conv_embeddings_file(conv_id)
        
        if os.path.exists(index_file):
            try:
                with open(index_file, "r", encoding="utf-8") as f:
                    self.documents = json.load(f)
                debug_log("RAG", f"Loaded {len(self.documents)} chunks")
                
                # Load or regenerate embeddings
                if os.path.exists(embeddings_file):
                    self._embeddings = np.load(embeddings_file)
                    debug_log("RAG", f"Loaded embeddings: {self._embeddings.shape}")
                elif self.embedding_model.is_loaded and self.documents:
                    debug_log("RAG", "Regenerating embeddings...")
                    texts = [c["content"] for c in self.documents]
                    self._embeddings = self.embedding_model.encode(texts)
                    if self._embeddings is not None:
                        np.save(embeddings_file, self._embeddings)
                        debug_log("RAG", f"✓ Embeddings: {self._embeddings.shape}")
                else:
                    self._embeddings = None
                return
            except Exception as e:
                debug_log("RAG", f"Error loading: {e}")
        
        self.documents = []
        self._embeddings = None
    
    def _save_index(self, log):
        if not self._current_conv_id:
            return
        try:
            index_file = self._get_conv_index_file(self._current_conv_id)
            with open(index_file, "w", encoding="utf-8") as f:
                json.dump(self.documents, f, ensure_ascii=False)
            
            if self._embeddings is not None:
                np.save(self._get_conv_embeddings_file(self._current_conv_id), self._embeddings)
            log(f"Saved {len(self.documents)} chunks")
        except Exception as e:
            log(f"Save error: {e}")
    
    def _split_text(self, text: str) -> List[str]:
        text = text.strip()
        if len(text) < 50:
            return [text] if len(text) > 20 else []
        
        # Split by sentences
        sentences = re.split(r'(?<=[.!?])\s+', text)
        chunks, current, length = [], [], 0
        
        for sent in sentences:
            words = sent.split()
            if length + len(words) <= config.RAG_CHUNK_SIZE:
                current.extend(words)
                length += len(words)
            else:
                if current:
                    chunks.append(" ".join(current))
                # Keep overlap
                overlap_start = max(0, len(current) - config.RAG_CHUNK_OVERLAP)
                current = current[overlap_start:] + words
                length = len(current)
        
        if current and len(current) >= 15:
            chunks.append(" ".join(current))
        
        return chunks
    
    def search(self, query: str) -> Tuple[str, List[dict]]:
        debug_log("RAG", f"Search: '{query[:40]}...' | docs={len(self.documents)} | emb={self._embeddings is not None}")
        
        if not self.documents:
            self.last_sources = []
            return "", []
        
        results = []
        
        # Try semantic search first
        if self._embeddings is not None and self.embedding_model.is_loaded:
            debug_log("RAG", "Using SEMANTIC search")
            query_emb = self.embedding_model.encode([query], is_query=True)
            
            if query_emb is not None:
                sims = np.dot(self._embeddings, query_emb[0])
                debug_log("RAG", f"Sims: min={sims.min():.3f} max={sims.max():.3f} mean={sims.mean():.3f}")
                
                top_indices = np.argsort(sims)[-config.RAG_TOP_K*2:][::-1]
                for idx in top_indices:
                    if sims[idx] >= config.RAG_MIN_SCORE:
                        results.append((self.documents[idx], float(sims[idx])))
                results = results[:config.RAG_TOP_K]
        
        # Fallback to keyword search
        if not results:
            debug_log("RAG", "Using KEYWORD search")
            query_words = set(query.lower().split())
            scored = []
            for doc in self.documents:
                content = doc["content"].lower()
                # Count word matches + partial matches
                score = sum(2 for w in query_words if w in content.split())
                score += sum(0.5 for w in query_words if w in content)
                if score > 0:
                    scored.append((doc, score))
            scored.sort(key=lambda x: x[1], reverse=True)
            results = scored[:config.RAG_TOP_K]
        
        if not results:
            self.last_sources = []
            return "", []
        
        self.last_sources = []
        context_parts = []
        
        for i, (doc, score) in enumerate(results, 1):
            self.last_sources.append({
                "index": i, 
                "source": doc["source"], 
                "chunk_id": doc["chunk_id"],
                "score": score, 
                "preview": doc["content"][:200]
            })
            context_parts.append(f"[Source: {doc['source']}]\n{doc['content']}")
            debug_log("RAG", f"  #{i}: {doc['source']} chunk={doc['chunk_id']} score={score:.3f}")
        
        return "\n\n".join(context_parts), self.last_sources
    
    def format_sources_for_display(self) -> str:
        if not self.last_sources:
            return ""
        lines = ["\n📚 Sources:"]
        for s in self.last_sources:
            lines.append(f"  [{s['index']}] {s['source']} (score: {s['score']:.2f})")
        return "\n".join(lines)
    
    def add_documents(self, file_paths: list, on_progress=None) -> bool:
        def log(msg):
            debug_log("RAG", msg)
            if on_progress:
                on_progress(msg)
        
        if not self._current_conv_id:
            log("No conversation!")
            return False
        
        docs_folder = self._get_conv_docs_folder(self._current_conv_id)
        new_chunks = []
        
        for fp in file_paths:
            fname = os.path.basename(fp)
            ext = os.path.splitext(fname)[1].lower()
            
            if ext not in DocumentParser.SUPPORTED_EXTENSIONS:
                log(f"⚠ {fname}: unsupported format")
                continue
            
            # Copy to conversation folder
            dest = os.path.join(docs_folder, fname)
            shutil.copy2(fp, dest)
            
            # Parse
            text = self.parser.parse(dest)
            if not text or len(text) < 10:
                log(f"⚠ {fname}: no text extracted")
                continue
            
            # Chunk
            chunks = self._split_text(text)
            for i, chunk in enumerate(chunks):
                new_chunks.append({
                    "source": fname, 
                    "chunk_id": i, 
                    "content": chunk
                })
            
            log(f"✓ {fname}: {len(chunks)} chunks ({len(text)} chars)")
        
        if not new_chunks:
            return False
        
        self.documents.extend(new_chunks)
        log(f"Total: {len(self.documents)} chunks")
        
        # Generate embeddings
        if self.embedding_model.is_loaded:
            log("Encoding embeddings...")
            texts = [c["content"] for c in self.documents]
            self._embeddings = self.embedding_model.encode(texts)
            if self._embeddings is not None:
                log(f"✓ Embeddings: {self._embeddings.shape}")
            else:
                log("⚠ Embedding failed, using keyword search")
        else:
            self._embeddings = None
            log("⚠ No embedding model, using keyword search")
        
        self._save_index(log)
        return True
    
    def clear_conversation_documents(self, conv_id: str):
        for f in [self._get_conv_index_file(conv_id), self._get_conv_embeddings_file(conv_id)]:
            if os.path.exists(f):
                try:
                    os.remove(f)
                except:
                    pass
        
        folder = os.path.join(self.user_docs_folder, conv_id)
        if os.path.exists(folder):
            try:
                shutil.rmtree(folder)
            except:
                pass
        
        if self._current_conv_id == conv_id:
            self.documents = []
            self._embeddings = None

# ============================================================================
# LLM ENGINE
# ============================================================================

class LLMEngine:
    def __init__(self):
        self.llm = None
        self.history = []
        self.rag = RAG()
        self.is_ready = False
        self._max_history = 3
    
    def load(self, on_progress=None) -> bool:
        def log(msg):
            debug_log("LLM", msg)
            if on_progress:
                on_progress(msg)
        
        try:
            # Initialize RAG first (loads embedding model)
            log("Initializing RAG...")
            self.rag.initialize(log)
            
            # Load LLM
            log("Importing llama_cpp...")
            from llama_cpp import Llama
            
            model_path = get_resource_path(config.MODEL_FILE)
            if not os.path.exists(model_path):
                raise FileNotFoundError(f"Model not found: {model_path}")
            
            file_size = os.path.getsize(model_path) / (1024*1024)
            log(f"Loading model: {model_path} ({file_size:.1f}MB)")
            
            self.llm = Llama(
                model_path=model_path,
                n_ctx=config.CONTEXT_SIZE,
                n_threads=config.THREADS,
                n_gpu_layers=config.GPU_LAYERS,
                verbose=False
            )
            
            self.is_ready = True
            log("✓ LLM Ready!")
            return True
            
        except Exception as e:
            log(f"✗ Error: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def _build_prompt(self, message: str, context: str = "") -> str:
        system = config.SYSTEM_PROMPT
        if context:
            system += f"\n\n=== RELEVANT DOCUMENTS ===\n{context}\n=== END DOCUMENTS ==="
        
        prompt = f"<|im_start|>system\n{system}<|im_end|>\n"
        
        for h in self.history[-self._max_history:]:
            prompt += f"<|im_start|>user\n{h['user']}<|im_end|>\n"
            prompt += f"<|im_start|>assistant\n{h['assistant']}<|im_end|>\n"
        
        prompt += f"<|im_start|>user\n{message}<|im_end|>\n<|im_start|>assistant\n"
        return prompt
    
    def generate(self, message: str) -> Iterator[str]:
        if not self.is_ready:
            yield "Model not ready"
            return
        
        # Search documents
        context, sources = "", []
        if config.RAG_ENABLED and self.rag.documents:
            context, sources = self.rag.search(message)
            debug_log("LLM", f"Context: {len(context)} chars, {len(sources)} sources")
        
        prompt = self._build_prompt(message, context)
        debug_log("LLM", f"Prompt: {len(prompt)} chars")
        
        full_response = ""
        
        try:
            for chunk in self.llm(
                prompt, 
                max_tokens=config.MAX_TOKENS, 
                stop=config.STOP_TOKENS,
                temperature=config.TEMPERATURE,
                top_p=config.TOP_P,
                repeat_penalty=config.REPEAT_PENALTY,
                stream=True
            ):
                token = chunk["choices"][0]["text"]
                full_response += token
                yield token
        except Exception as e:
            debug_log("LLM", f"Generation error: {e}")
            yield f"\n[Error: {e}]"
        
        # Add sources
        if sources and config.RAG_SHOW_SOURCES:
            src_text = self.rag.format_sources_for_display()
            yield src_text
        
        # Save to history
        if full_response.strip():
            clean = full_response.split("📚")[0].strip()
            self.history.append({"user": message, "assistant": clean})
            if len(self.history) > self._max_history * 2:
                self.history = self.history[-self._max_history:]
    
    def clear_history(self):
        self.history = []

# ============================================================================
# CONVERSATION MANAGER
# ============================================================================

@dataclass
class Conversation:
    id: str
    title: str
    created_at: str
    updated_at: str
    messages: List[Dict[str, str]]
    document_ids: List[str]
    
    def to_dict(self): 
        return asdict(self)
    
    @classmethod
    def from_dict(cls, d): 
        return cls(**d)

class ConversationManager:
    def __init__(self):
        self.conversations = {}
        self.current_id = None
        self.data_dir = get_writable_path("conversations")
        self.index_file = os.path.join(self.data_dir, "index.json")
        os.makedirs(self.data_dir, exist_ok=True)
        self._load_index()
    
    def _load_index(self):
        if os.path.exists(self.index_file):
            try:
                with open(self.index_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                self.conversations = {
                    k: Conversation.from_dict(v) 
                    for k, v in data.get("conversations", {}).items()
                }
                self.current_id = data.get("current_id")
            except Exception as e:
                debug_log("CONV", f"Load error: {e}")
    
    def _save_index(self):
        try:
            data = {
                "conversations": {k: v.to_dict() for k, v in self.conversations.items()}, 
                "current_id": self.current_id
            }
            with open(self.index_file, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            debug_log("CONV", f"Save error: {e}")
    
    def create_conversation(self, title="") -> Conversation:
        cid = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        now = datetime.now().isoformat()
        conv = Conversation(
            cid, 
            title or f"Chat {datetime.now().strftime('%d/%m %H:%M')}", 
            now, now, [], []
        )
        self.conversations[cid] = conv
        self.current_id = cid
        self._save_index()
        return conv
    
    def get_current(self): 
        return self.conversations.get(self.current_id)
    
    def set_current(self, cid): 
        if cid in self.conversations:
            self.current_id = cid
            self._save_index()
            return self.conversations[cid]
        return None
    
    def get_all(self): 
        return sorted(self.conversations.values(), key=lambda c: c.updated_at, reverse=True)
    
    def add_message(self, role, content):
        conv = self.get_current() or self.create_conversation()
        conv.messages.append({"role": role, "content": content})
        conv.updated_at = datetime.now().isoformat()
        if role == "user" and len(conv.messages) == 1:
            conv.title = content[:40] + ("..." if len(content) > 40 else "")
        self._save_index()
    
    def add_document(self, fname):
        conv = self.get_current()
        if conv and fname not in conv.document_ids:
            conv.document_ids.append(fname)
            self._save_index()
    
    def delete_conversation(self, cid):
        if cid in self.conversations:
            del self.conversations[cid]
            if self.current_id == cid:
                remaining = self.get_all()
                self.current_id = remaining[0].id if remaining else None
            self._save_index()
    
    def clear_history(self):
        conv = self.get_current()
        if conv:
            conv.messages = []
            self._save_index()

# ============================================================================
# UI
# ============================================================================

ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")

class ConversationItem(ctk.CTkFrame):
    def __init__(self, parent, conv_id, title, doc_count, is_active, on_select, on_delete):
        super().__init__(parent, corner_radius=8, height=50)
        self.conv_id = conv_id
        
        color = ("#3b7ac7", "#2d5f9e") if is_active else ("gray75", "gray25")
        self.configure(fg_color=color)
        self.bind("<Button-1>", lambda e: on_select(conv_id))
        
        display_title = title[:25] + "..." if len(title) > 25 else title
        lbl = ctk.CTkLabel(
            self, 
            text=display_title,
            font=ctk.CTkFont(size=12, weight="bold" if is_active else "normal")
        )
        lbl.pack(side="left", padx=10, fill="x", expand=True)
        lbl.bind("<Button-1>", lambda e: on_select(conv_id))
        
        if doc_count:
            ctk.CTkLabel(
                self, 
                text=f"📄{doc_count}", 
                font=ctk.CTkFont(size=10),
                text_color=("#50fa7b", "#40c969")
            ).pack(side="left", padx=5)
        
        ctk.CTkButton(
            self, 
            text="✕", 
            width=24, 
            height=24, 
            fg_color="transparent",
            hover_color=("#ff5555", "#cc4444"), 
            command=lambda: on_delete(conv_id)
        ).pack(side="right", padx=5)

class Sidebar(ctk.CTkFrame):
    def __init__(self, parent, on_new, on_select, on_delete):
        super().__init__(parent, width=250, corner_radius=0)
        self.on_new = on_new
        self.on_select = on_select
        self.on_delete = on_delete
        self.pack_propagate(False)
        
        # Header
        header = ctk.CTkFrame(self, fg_color="transparent", height=60)
        header.pack(fill="x", padx=10, pady=10)
        
        ctk.CTkLabel(
            header, 
            text="💬 Chats", 
            font=ctk.CTkFont(size=18, weight="bold")
        ).pack(side="left")
        
        ctk.CTkButton(
            header, 
            text="+ New", 
            width=70, 
            height=32, 
            fg_color=("#50fa7b", "#40c969"),
            text_color="black", 
            command=on_new
        ).pack(side="right")
        
        # Conversation list
        self.conv_list = ctk.CTkScrollableFrame(self, fg_color="transparent")
        self.conv_list.pack(fill="both", expand=True, padx=5, pady=5)
    
    def update_conversations(self, convs, current_id):
        for w in self.conv_list.winfo_children():
            w.destroy()
        
        for conv in convs:
            def make_delete_handler(cid):
                def handler():
                    if messagebox.askyesno("Delete?", "Delete this chat?"):
                        self.on_delete(cid)
                return handler
            
            item = ConversationItem(
                self.conv_list, 
                conv.id, 
                conv.title, 
                len(conv.document_ids),
                conv.id == current_id, 
                self.on_select, 
                make_delete_handler(conv.id)
            )
            item.pack(fill="x", pady=2)

class ChatUI:
    def __init__(self, root, on_send, on_clear, on_load, on_new, on_select, on_delete):
        self.root = root
        self.on_send = on_send
        self.on_clear = on_clear
        self.on_load = on_load
        
        root.title(f"🤖 {config.APP_NAME}")
        root.geometry(config.WINDOW_SIZE)
        root.grid_rowconfigure(0, weight=1)
        root.grid_columnconfigure(1, weight=1)
        
        # Sidebar
        self.sidebar = Sidebar(root, on_new, on_select, on_delete)
        self.sidebar.grid(row=0, column=0, sticky="nsew")
        
        # Main area
        main = ctk.CTkFrame(root, fg_color="transparent")
        main.grid(row=0, column=1, sticky="nsew", padx=10, pady=10)
        main.grid_rowconfigure(2, weight=1)
        main.grid_columnconfigure(0, weight=1)
        
        # Header
        header = ctk.CTkFrame(main, fg_color="transparent")
        header.grid(row=0, column=0, sticky="ew")
        
        ctk.CTkLabel(
            header, 
            text=f"🤖 {config.APP_NAME}", 
            font=ctk.CTkFont(size=24, weight="bold")
        ).pack(side="left")
        
        self.status = ctk.CTkLabel(header, text="⏳ Loading...", text_color="gray")
        self.status.pack(side="right")
        
        # Toolbar
        toolbar = ctk.CTkFrame(main, fg_color="transparent")
        toolbar.grid(row=1, column=0, sticky="ew", pady=5)
        
        ctk.CTkButton(
            toolbar, 
            text="📁 Load Files", 
            command=self._load_files, 
            width=100
        ).pack(side="left", padx=5)
        
        ctk.CTkButton(
            toolbar, 
            text="🗑️ Clear", 
            command=on_clear, 
            width=80, 
            fg_color="#ff5555"
        ).pack(side="left")
        
        self.doc_info = ctk.CTkLabel(toolbar, text="📚 No docs", text_color="gray")
        self.doc_info.pack(side="right")
        
        # Chat area
        self.chat = ctk.CTkTextbox(
            main, 
            wrap="word", 
            font=ctk.CTkFont(family="Consolas", size=12)
        )
        self.chat.grid(row=2, column=0, sticky="nsew", pady=5)
        
        # Input area
        inp_frame = ctk.CTkFrame(main, fg_color="transparent")
        inp_frame.grid(row=3, column=0, sticky="ew")
        inp_frame.grid_columnconfigure(0, weight=1)
        
        self.input = ctk.CTkTextbox(
            inp_frame, 
            height=60, 
            font=ctk.CTkFont(family="Consolas", size=12)
        )
        self.input.grid(row=0, column=0, sticky="ew", padx=(0, 10))
        self.input.bind("<Return>", self._on_enter)
        
        ctk.CTkButton(
            inp_frame, 
            text="Send ➤", 
            command=self._send, 
            width=90, 
            height=60,
            fg_color=("#50fa7b", "#40c969"), 
            text_color="black"
        ).grid(row=0, column=1)
    
    def _on_enter(self, e):
        if not (e.state & 1):  # Not shift+enter
            self._send()
            return "break"
    
    def _send(self):
        text = self.input.get("0.0", "end").strip()
        if text:
            self.input.delete("0.0", "end")
            self.on_send(text)
    
    def _load_files(self):
        files = filedialog.askopenfilenames(
            title="Select documents", 
            filetypes=[
                ("All supported", "*.txt *.md *.pdf *.xlsx *.docx *.pptx *.json *.csv *.xml *.yaml"),
                ("Documents", "*.pdf *.docx *.txt *.md"),
                ("Spreadsheets", "*.xlsx *.csv"),
                ("All", "*.*")
            ]
        )
        if files:
            self.on_load(files)
    
    def set_status(self, text, error=False):
        color = "#ff5555" if error else "#50fa7b"
        self.status.configure(text=text, text_color=color)
    
    def update_doc_count(self, count, chunks=0, semantic=False):
        if count == 0:
            self.doc_info.configure(text="📚 No docs", text_color="gray")
        else:
            text = f"📚 {count} doc{'s' if count > 1 else ''} ({chunks} chunks)"
            if semantic:
                text += " 🧠"
            self.doc_info.configure(text=text, text_color="#50fa7b")
    
    def add_message(self, sender, text, tag=""):
        if self.chat.get("0.0", "end").strip():
            self.chat.insert("end", "\n")
        self.chat.insert("end", f"{sender}:\n{text}\n")
        self.chat.see("end")
    
    def stream(self, text):
        self.chat.insert("end", text)
        self.chat.see("end")
    
    def clear_chat(self):
        self.chat.delete("0.0", "end")
    
    def set_enabled(self, enabled):
        self.input.configure(state="normal" if enabled else "disabled")
    
    def update_sidebar(self, convs, current_id):
        self.sidebar.update_conversations(convs, current_id)
    
    def load_messages(self, msgs):
        self.clear_chat()
        for m in msgs:
            sender = "You" if m["role"] == "user" else "Assistant"
            self.add_message(sender, m["content"])

# ============================================================================
# APP
# ============================================================================

class App:
    def __init__(self):
        self.root = ctk.CTk()
        self.engine = LLMEngine()
        self.conv_manager = ConversationManager()
        self.generating = False
        
        self.ui = ChatUI(
            self.root, 
            self.send, 
            self.clear, 
            self.load_files,
            self.new_chat, 
            self.select_chat, 
            self.delete_chat
        )
        
        self.ui.add_message("System", "🚀 Starting Local Chat v4...", "system")
        threading.Thread(target=self._load, daemon=True).start()
    
    def _load(self):
        success = self.engine.load(
            lambda m: self.root.after(0, lambda: self.ui.add_message("Debug", m))
        )
        if success:
            self.root.after(0, self._ready)
        else:
            self.root.after(0, lambda: self.ui.set_status("❌ Load Failed", True))
    
    def _ready(self):
        self.ui.set_status("✅ Ready!")
        
        if not self.conv_manager.get_all():
            self.conv_manager.create_conversation("Welcome")
        
        self._load_conv()
        self._update_sidebar()
        
        emb_status = "🧠 Semantic" if self.engine.rag.embedding_model.is_loaded else "⚠ Keyword"
        self.ui.add_message("System", f"✨ Ready! {emb_status} search enabled.")
    
    def _update_sidebar(self):
        convs = self.conv_manager.get_all()
        current = self.conv_manager.get_current()
        self.ui.update_sidebar(convs, current.id if current else None)
    
    def _load_conv(self):
        conv = self.conv_manager.get_current()
        if not conv:
            self.ui.clear_chat()
            self.ui.update_doc_count(0)
            return
        
        self.engine.rag.load_conversation_documents(conv.id)
        self.ui.load_messages(conv.messages)
        
        has_emb = self.engine.rag._embeddings is not None
        self.ui.update_doc_count(len(conv.document_ids), len(self.engine.rag.documents), has_emb)
        self.engine.clear_history()
    
    def new_chat(self):
        self.conv_manager.create_conversation()
        self._load_conv()
        self._update_sidebar()
        self.ui.add_message("System", "💬 New chat created!")
    
    def select_chat(self, cid):
        self.conv_manager.set_current(cid)
        self._load_conv()
        self._update_sidebar()
    
    def delete_chat(self, cid):
        self.engine.rag.clear_conversation_documents(cid)
        self.conv_manager.delete_conversation(cid)
        self._load_conv()
        self._update_sidebar()
    
    def send(self, msg):
        if self.generating or not self.engine.is_ready:
            return
        
        if not self.conv_manager.get_current():
            self.conv_manager.create_conversation()
            self._update_sidebar()
        
        self.ui.add_message("You", msg)
        self.conv_manager.add_message("user", msg)
        
        self.generating = True
        self.ui.set_enabled(False)
        self.ui.set_status("🤔 Thinking...")
        
        threading.Thread(target=self._generate, args=(msg,), daemon=True).start()
    
    def _generate(self, msg):
        self.root.after(0, lambda: self.ui.add_message("Assistant", ""))
        
        full_response = ""
        for token in self.engine.generate(msg):
            full_response += token
            self.root.after(0, lambda t=token: self.ui.stream(t))
        
        # Save clean response (without sources)
        clean = full_response.split("📚")[0].strip()
        if clean:
            self.conv_manager.add_message("assistant", clean)
        
        self.root.after(0, self._done)
    
    def _done(self):
        self.generating = False
        self.ui.set_enabled(True)
        self.ui.set_status("✅ Ready!")
        self._update_sidebar()
    
    def clear(self):
        self.conv_manager.clear_history()
        self.engine.clear_history()
        self.ui.clear_chat()
        self.ui.add_message("System", "🗑️ Chat cleared.")
    
    def load_files(self, files):
        conv = self.conv_manager.get_current()
        if not conv:
            self.conv_manager.create_conversation()
            conv = self.conv_manager.get_current()
            self._update_sidebar()
        
        self.engine.rag._current_conv_id = conv.id
        self.ui.add_message("System", f"📥 Loading {len(files)} file(s)...")
        self.ui.set_status("⏳ Loading files...")
        
        def task():
            success = self.engine.rag.add_documents(
                files, 
                lambda m: self.root.after(0, lambda: self.ui.add_message("Debug", m))
            )
            
            if success:
                for f in files:
                    self.conv_manager.add_document(os.path.basename(f))
                
                conv = self.conv_manager.get_current()
                has_emb = self.engine.rag._embeddings is not None
                
                self.root.after(0, lambda: self.ui.update_doc_count(
                    len(conv.document_ids) if conv else 0, 
                    len(self.engine.rag.documents), 
                    has_emb
                ))
                
                search_type = "🧠 Semantic" if has_emb else "⚠ Keyword"
                self.root.after(0, lambda: self.ui.add_message(
                    "System", 
                    f"✅ Documents loaded! {search_type} search ready."
                ))
                self.root.after(0, lambda: self.ui.set_status("✅ Ready!"))
                self.root.after(0, self._update_sidebar)
            else:
                self.root.after(0, lambda: self.ui.add_message("Error", "❌ Failed to load documents"))
                self.root.after(0, lambda: self.ui.set_status("❌ Load failed", True))
        
        threading.Thread(target=task, daemon=True).start()
    
    def run(self):
        self.root.mainloop()

# ============================================================================
# ENTRY POINT
# ============================================================================

if __name__ == "__main__":
    debug_log("MAIN", f"Python {sys.version}")
    debug_log("MAIN", f"Platform: {sys.platform}")
    debug_log("MAIN", f"Executable: {sys.executable}")
    debug_log("MAIN", f"CWD: {os.getcwd()}")
    if hasattr(sys, '_MEIPASS'):
        debug_log("MAIN", f"PyInstaller _MEIPASS: {sys._MEIPASS}")
    App().run()
