"""
OCR processor for scanned documents and images.
"""

import os
import sys
import glob
from typing import Optional, Callable

# OCR dependencies
try:
    import pytesseract
    from PIL import Image
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False


class OCRProcessor:
    """Handles OCR for scanned documents and images."""
    
    def __init__(self):
        self.available = HAS_OCR
        self.pdf_support = HAS_PDF2IMAGE
        self.poppler_path = None
        self._configure_tesseract()
        self._configure_poppler()
    
    def _configure_tesseract(self):
        """Configure Tesseract path for PyInstaller on Windows."""
        if not HAS_OCR:
            return
        
        if sys.platform == 'win32':
            # Try bundled path first (if we bundle it in future)
            if hasattr(sys, '_MEIPASS'):
                tesseract_path = os.path.join(sys._MEIPASS, 'tesseract', 'tesseract.exe')
                if os.path.exists(tesseract_path):
                    pytesseract.pytesseract.tesseract_cmd = tesseract_path
                    print(f"[OCR] Using bundled Tesseract: {tesseract_path}")
                    return
            
            # Try standard Windows installation paths
            standard_paths = [
                r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
                os.path.join(os.environ.get('LOCALAPPDATA', ''), 'Programs', 'Tesseract-OCR', 'tesseract.exe'),
            ]
            
            for path in standard_paths:
                if os.path.exists(path):
                    pytesseract.pytesseract.tesseract_cmd = path
                    print(f"[OCR] Using Tesseract: {path}")
                    return
            
            print("[OCR] Warning: Tesseract not found in standard paths")
        
        # Verify Tesseract is accessible
        try:
            version = pytesseract.get_tesseract_version()
            print(f"[OCR] Tesseract version: {version}")
        except Exception as e:
            print(f"[OCR] Tesseract not accessible: {e}")
            self.available = False
    
    def _configure_poppler(self):
        """Configure Poppler path for pdf2image on Windows."""
        if not HAS_PDF2IMAGE:
            return
        
        if sys.platform == 'win32':
            # Common Poppler installation paths on Windows
            poppler_search_paths = [
                r"C:\Program Files\poppler-*\Library\bin",
                r"C:\Program Files\poppler\Library\bin",
                r"C:\Program Files (x86)\poppler-*\Library\bin",
                r"C:\ProgramData\chocolatey\lib\poppler\tools\Library\bin",
                r"C:\ProgramData\chocolatey\bin",
                os.path.join(os.environ.get('LOCALAPPDATA', ''), 'poppler-*', 'Library', 'bin'),
            ]
            
            for pattern in poppler_search_paths:
                # Handle glob patterns
                if '*' in pattern:
                    matches = glob.glob(pattern)
                    for match in matches:
                        if os.path.isdir(match):
                            self.poppler_path = match
                            print(f"[OCR] Using Poppler: {self.poppler_path}")
                            return
                elif os.path.isdir(pattern):
                    self.poppler_path = pattern
                    print(f"[OCR] Using Poppler: {self.poppler_path}")
                    return
            
            # Check if pdftoppm is in PATH
            import shutil
            if shutil.which('pdftoppm'):
                print("[OCR] Poppler found in PATH")
                self.poppler_path = None  # Will use PATH
            else:
                print("[OCR] Warning: Poppler not found, PDF OCR may not work")
    
    def get_status(self) -> dict:
        """Returns OCR availability status."""
        status = {
            "ocr_available": self.available,
            "pdf_ocr_available": self.available and self.pdf_support,
        }
        
        # Get available languages
        if self.available:
            try:
                langs = pytesseract.get_languages()
                status["available_languages"] = langs
            except Exception:
                status["available_languages"] = []
        
        return status
    
    def _get_best_language(self, preferred: str = 'eng') -> str:
        """Get the best available language for OCR."""
        if not self.available:
            return 'eng'
        
        try:
            available_langs = pytesseract.get_languages()
            
            # Parse preferred languages (e.g., 'fra+eng')
            preferred_list = preferred.split('+')
            
            # Filter to only available languages
            valid_langs = [lang for lang in preferred_list if lang in available_langs]
            
            if valid_langs:
                return '+'.join(valid_langs)
            
            # Fallback to English if available
            if 'eng' in available_langs:
                return 'eng'
            
            # Use first available language
            if available_langs:
                return available_langs[0]
            
            return 'eng'
        except Exception:
            return 'eng'
    
    def ocr_image(self, image_path: str, lang: str = 'eng') -> str:
        """Extract text from an image using OCR.
        
        Args:
            image_path: Path to the image file
            lang: Tesseract language codes (default: English)
        
        Returns:
            Extracted text or empty string on failure
        """
        if not self.available:
            print("[OCR] OCR not available")
            return ""
        
        try:
            image = Image.open(image_path)
            
            # Convert to RGB if necessary (handles RGBA, grayscale, etc.)
            if image.mode not in ('RGB', 'L'):
                image = image.convert('RGB')
            
            # Get best available language
            use_lang = self._get_best_language(lang)
            
            try:
                text = pytesseract.image_to_string(image, lang=use_lang)
            except pytesseract.TesseractError as e:
                print(f"[OCR] Language error with '{use_lang}': {e}")
                # Fallback to no language specification
                text = pytesseract.image_to_string(image)
            
            return text.strip()
        except Exception as e:
            print(f"[OCR] Error processing image {image_path}: {e}")
            return ""
    
    def ocr_image_from_bytes(self, image_bytes: bytes, lang: str = 'eng') -> str:
        """Extract text from image bytes using OCR.
        
        Args:
            image_bytes: Raw image bytes
            lang: Tesseract language codes
        
        Returns:
            Extracted text or empty string on failure
        """
        if not self.available:
            return ""
        
        try:
            from io import BytesIO
            image = Image.open(BytesIO(image_bytes))
            
            if image.mode not in ('RGB', 'L'):
                image = image.convert('RGB')
            
            # Get best available language
            use_lang = self._get_best_language(lang)
            
            try:
                text = pytesseract.image_to_string(image, lang=use_lang)
            except pytesseract.TesseractError as e:
                print(f"[OCR] Language error: {e}")
                text = pytesseract.image_to_string(image)
            
            return text.strip()
        except Exception as e:
            print(f"[OCR] Error processing image bytes: {e}")
            return ""
    
    def ocr_pdf(self, pdf_path: str, lang: str = 'eng', 
                dpi: int = 300, on_progress: Optional[Callable[[str], None]] = None) -> str:
        """Extract text from a scanned PDF using OCR.
        
        Args:
            pdf_path: Path to the PDF file
            lang: Tesseract language codes
            dpi: Resolution for PDF to image conversion
            on_progress: Optional callback for progress updates
        
        Returns:
            Extracted text or empty string on failure
        """
        if not self.available:
            print("[OCR] OCR not available")
            return ""
        
        if not self.pdf_support:
            print("[OCR] pdf2image not available")
            return ""
        
        def log(msg):
            print(f"[OCR] {msg}")
            if on_progress:
                on_progress(msg)
        
        try:
            log(f"Converting PDF to images (dpi={dpi})...")
            
            # Build convert_from_path arguments
            convert_kwargs = {'dpi': dpi}
            
            # Add poppler path on Windows if found
            if sys.platform == 'win32' and self.poppler_path:
                convert_kwargs['poppler_path'] = self.poppler_path
            
            try:
                images = convert_from_path(pdf_path, **convert_kwargs)
            except Exception as e:
                log(f"PDF conversion error: {e}")
                # Try with lower DPI
                if dpi > 150:
                    log("Retrying with lower DPI (150)...")
                    convert_kwargs['dpi'] = 150
                    images = convert_from_path(pdf_path, **convert_kwargs)
                else:
                    raise
            
            all_text = []
            total_pages = len(images)
            
            # Get best available language
            use_lang = self._get_best_language(lang)
            log(f"Using language: {use_lang}")
            
            for i, image in enumerate(images):
                log(f"OCR page {i + 1}/{total_pages}")
                
                # Convert to RGB if needed
                if image.mode not in ('RGB', 'L'):
                    image = image.convert('RGB')
                
                try:
                    text = pytesseract.image_to_string(image, lang=use_lang)
                except pytesseract.TesseractError:
                    text = pytesseract.image_to_string(image)
                
                if text.strip():
                    all_text.append(f"=== Page {i + 1} ===\n{text.strip()}")
            
            log(f"OCR complete: {len(all_text)} pages with text")
            return "\n\n".join(all_text)
        except Exception as e:
            log(f"Error processing PDF: {e}")
            import traceback
            traceback.print_exc()
            return ""
    
    def ocr_pptx_images(self, pptx_path: str, lang: str = 'eng') -> str:
        """Extract text from images embedded in PowerPoint.
        
        Args:
            pptx_path: Path to the PPTX file
            lang: Tesseract language codes
        
        Returns:
            Extracted text from all images or empty string
        """
        if not self.available:
            return ""
        
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            ocr_texts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image_bytes = shape.image.blob
                            text = self.ocr_image_from_bytes(image_bytes, lang)
                            if text:
                                ocr_texts.append(f"[Slide {slide_num} - Image]\n{text}")
                        except Exception as e:
                            print(f"[OCR] Error with image on slide {slide_num}: {e}")
            
            return "\n\n".join(ocr_texts)
        except ImportError:
            print("[OCR] python-pptx not available")
            return ""
        except Exception as e:
            print(f"[OCR] Error processing PPTX images: {e}")
            return ""
    
    def ocr_docx_images(self, docx_path: str, lang: str = 'eng') -> str:
        """Extract text from images embedded in Word documents.
        
        Args:
            docx_path: Path to the DOCX file
            lang: Tesseract language codes
        
        Returns:
            Extracted text from all images or empty string
        """
        if not self.available:
            return ""
        
        try:
            from docx import Document
            
            doc = Document(docx_path)
            ocr_texts = []
            image_count = 0
            
            # Access images through document parts
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_bytes = rel.target_part.blob
                        text = self.ocr_image_from_bytes(image_bytes, lang)
                        if text:
                            image_count += 1
                            ocr_texts.append(f"[Image {image_count}]\n{text}")
                    except Exception as e:
                        print(f"[OCR] Error with image in DOCX: {e}")
            
            return "\n\n".join(ocr_texts)
        except ImportError:
            print("[OCR] python-docx not available")
            return ""
        except Exception as e:
            print(f"[OCR] Error processing DOCX images: {e}")
            return ""


# Convenience function for quick OCR check
def is_ocr_available() -> bool:
    """Quick check if OCR is available."""
    processor = OCRProcessor()
    return processor.available
